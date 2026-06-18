# -*- coding: utf-8 -*-
"""음악 데이팅 앱 → 운영체제(OS) 개념 학습용 PPT 생성기."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- 브랜드/팔레트 ----
CORAL = RGBColor(0xFF, 0x62, 0x5D)
CORAL_SOFT = RGBColor(0xFF, 0xEE, 0xED)
INK = RGBColor(0x22, 0x22, 0x28)
GRAY = RGBColor(0x5B, 0x5B, 0x66)
LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x28)
CODE_FG = RGBColor(0xE6, 0xE6, 0xEC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            para = [para]
        for (t, sz, bold, col) in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.bold = bold
            r.font.color.rgb = col
            r.font.name = "맑은 고딕"
    return tb


def bg(s, color=WHITE):
    rect(s, 0, 0, SW, SH, color)


def header(s, kicker, title):
    rect(s, 0, 0, SW, Inches(1.35), CORAL_SOFT)
    rect(s, 0, 0, Inches(0.18), Inches(1.35), CORAL)
    txt(s, Inches(0.6), Inches(0.22), Inches(11), Inches(0.4),
        [[(kicker, 14, True, CORAL)]])
    txt(s, Inches(0.6), Inches(0.52), Inches(12), Inches(0.7),
        [[(title, 30, True, INK)]])


def code_box(s, x, y, w, h, lines, fontsize=12):
    sp = rect(s, x, y, w, h, CODE_BG)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.0
        p.space_after = Pt(1)
        col = CODE_FG
        if ln.strip().startswith("--") or ln.strip().startswith("//"):
            col = RGBColor(0x7A, 0x9E, 0x7E)
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(fontsize)
        r.font.name = "Consolas"
        r.font.color.rgb = col
    return sp


def card(s, x, y, w, h, title, body_lines, accent=CORAL):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, w, Inches(0.09), accent)
    txt(s, x + Inches(0.22), y + Inches(0.22), w - Inches(0.44), Inches(0.5),
        [[(title, 16, True, INK)]])
    runs = [[(b, 12.5, False, GRAY)] for b in body_lines]
    txt(s, x + Inches(0.22), y + Inches(0.78), w - Inches(0.44), h - Inches(0.9),
        runs, space_after=5, line_spacing=1.08)


def footer(s, n):
    txt(s, Inches(11.6), Inches(7.05), Inches(1.5), Inches(0.35),
        [[("뮤직매치 · OS 사례  %02d" % n, 9, False, GRAY)]], align=PP_ALIGN.RIGHT)


# =========================================================
# 1. 표지
# =========================================================
s = slide()
bg(s, INK)
rect(s, 0, 0, Inches(0.25), SH, CORAL)
rect(s, Inches(0.9), Inches(2.1), Inches(1.3), Inches(0.12), CORAL)
txt(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
    [[("내가 만든 데이팅 앱으로 배우는", 20, False, RGBColor(0xC9, 0xC9, 0xD2))],
     [("운영체제(OS) 핵심 개념", 46, True, WHITE)]], space_after=10)
txt(s, Inches(0.9), Inches(4.1), Inches(11), Inches(1.0),
    [[("음악 취향 기반 1:1 매칭 + 채팅 웹앱 (뮤직매치)", 17, False, RGBColor(0xE6, 0xE6, 0xEC))],
     [("Remix · Supabase(PostgreSQL) · Railway", 14, False, GRAY)]], space_after=6)
rect(s, Inches(0.9), Inches(5.6), Inches(5.6), Inches(0.9), RGBColor(0x2C, 0x2C, 0x38))
txt(s, Inches(1.1), Inches(5.72), Inches(5.3), Inches(0.7),
    [[("실제 코드 → 교과서 OS 개념으로 1:1 매핑", 13, True, CORAL)],
     [("동시성 · 임계구역 · 캐시 · 스케줄링 · 보호 · IPC", 11.5, False, RGBColor(0xC9, 0xC9, 0xD2))]],
    space_after=4)

# =========================================================
# 2. 목차 / 매핑 한눈에
# =========================================================
s = slide()
bg(s)
header(s, "OVERVIEW", "한눈에 보기 — 앱 기능 ↔ OS 개념 매핑")
rows = [
    ("동시 매칭 버그 수정 (advisory lock)", "임계구역·상호배제·경쟁 상태·락", "03~05"),
    ("멜론 검색/차트 인메모리 캐시 (TTL)", "캐싱·지역성·캐시 교체·만료", "06~07"),
    ("매칭 대기열 / 승인 대기 / 한 명 더", "프로세스 상태·스케줄링·큐", "08"),
    ("실시간 채팅 (Realtime 구독)", "프로세스 간 통신(IPC)·메시지 패싱", "09"),
    ("한 번만 보기 사진 / 펑 메시지", "자원 회수·파일 수명·동기화", "10"),
    ("RLS · service_role · 서버 전용 키", "보호·접근제어·이중 모드(권한)", "11"),
]
y = Inches(1.65)
# 헤더행
rect(s, Inches(0.6), y, Inches(12.1), Inches(0.5), CORAL)
txt(s, Inches(0.8), y + Inches(0.08), Inches(5.0), Inches(0.4), [[("앱에서 만든 기능", 13, True, WHITE)]])
txt(s, Inches(6.0), y + Inches(0.08), Inches(4.8), Inches(0.4), [[("대응되는 OS 개념", 13, True, WHITE)]])
txt(s, Inches(11.3), y + Inches(0.08), Inches(1.3), Inches(0.4), [[("슬라이드", 13, True, WHITE)]])
y = y + Inches(0.5)
for i, (a, b, p) in enumerate(rows):
    c = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.62), c)
    txt(s, Inches(0.8), y + Inches(0.13), Inches(5.1), Inches(0.4), [[(a, 13, False, INK)]])
    txt(s, Inches(6.0), y + Inches(0.13), Inches(5.0), Inches(0.4), [[(b, 13, True, CORAL)]])
    txt(s, Inches(11.3), y + Inches(0.13), Inches(1.3), Inches(0.4), [[(p, 13, False, GRAY)]])
    y = y + Inches(0.62)
footer(s, 2)

# =========================================================
# 3. 경쟁 상태(문제 정의)
# =========================================================
s = slide()
bg(s)
header(s, "동시성 #1 · 문제", "경쟁 상태(Race Condition) — 두 명이 동시에 '매칭'을 누르면?")
txt(s, Inches(0.6), Inches(1.6), Inches(12), Inches(0.7),
    [[("매칭 함수는 ① 후보를 SELECT 하고 ② matches 테이블에 INSERT 합니다. "
       "이 둘 사이에 다른 호출이 끼어들면(인터리빙) 잘못된 결과가 생깁니다.", 14, False, GRAY)]])
card(s, Inches(0.6), Inches(2.5), Inches(5.95), Inches(2.0), "(a) 빠른 더블클릭 / 두 탭",
     ["같은 사용자가 거의 동시에 두 번 호출", "둘 다 '나는 아직 매칭 없음'으로 판단",
      "→ 내가 매칭 2개에 동시에 들어감"])
card(s, Inches(6.75), Inches(2.5), Inches(5.95), Inches(2.0), "(b) A·B가 서로를 동시에 선택",
     ["A는 B를, B는 A를 동시에 후보로 집음", "각자 INSERT 실행",
      "→ 같은 (A,B) 쌍에 active 매칭 2행 생성"], accent=INK)
rect(s, Inches(0.6), Inches(4.75), Inches(12.1), Inches(1.9), CORAL_SOFT)
txt(s, Inches(0.85), Inches(4.9), Inches(11.6), Inches(0.5),
    [[("📖 OS 용어로 말하면", 15, True, CORAL)]])
txt(s, Inches(0.85), Inches(5.35), Inches(11.6), Inches(1.2),
    [[("• 임계구역(Critical Section): SELECT~INSERT 구간 — 한 번에 하나만 실행돼야 한다.", 13.5, False, INK)],
     [("• 경쟁 상태: 실행 순서(스케줄링)에 따라 결과가 달라지는 버그.", 13.5, False, INK)],
     [("• 원인: 공유 자원(matches 테이블)에 대한 상호배제(Mutual Exclusion) 부재.", 13.5, False, INK)]],
    space_after=6)
footer(s, 3)

# =========================================================
# 4. 해결(락) — 코드
# =========================================================
s = slide()
bg(s)
header(s, "동시성 #2 · 해결", "상호배제 — Advisory Lock 으로 임계구역 직렬화")
code_box(s, Inches(0.6), Inches(1.6), Inches(7.4), Inches(4.9), [
    "-- 빠른 경로: 이미 매칭 있으면 락 없이 즉시 반환 (멱등)",
    "select id into v_match_id from matches",
    "where status='active'",
    "  and (user_a=p_user_id or user_b=p_user_id);",
    "if v_match_id is not null then return v_match_id; end if;",
    "",
    "-- ① 임계구역 진입: 트랜잭션 advisory lock",
    "--    (= OS의 mutex.lock(). 같은 키는 한 번에 하나만)",
    "perform pg_advisory_xact_lock(hashtext('match_creation'));",
    "",
    "-- ② 락 획득 후 재확인 (double-checked):",
    "--    대기하는 동안 상대가 나를 매칭했을 수 있음",
    "select id into v_match_id from matches ...;",
    "if v_match_id is not null then return v_match_id; end if;",
    "",
    "-- ③ 후보 선정 → INSERT (이제 안전)",
    "insert into matches(user_a, user_b, status) ...;",
    "-- 트랜잭션 종료 시 락 자동 해제 (= mutex.unlock())",
], fontsize=12)
card(s, Inches(8.25), Inches(1.6), Inches(4.45), Inches(2.35), "핵심 3중 방어",
     ["① 락으로 구간 직렬화 (한 명씩)",
      "② 락 후 재확인 — 이중 검사 잠금",
      "③ DB 부분 유니크 인덱스로 최후 차단"], accent=CORAL)
card(s, Inches(8.25), Inches(4.15), Inches(4.45), Inches(2.35), "OS 매핑",
     ["pg_advisory_xact_lock ≈ mutex/세마포어",
      "트랜잭션 종료 = 락 자동 반납",
      "'빠른 경로'는 무경합 시 락 비용 회피",
      "→ 동시성과 성능의 trade-off"], accent=INK)
footer(s, 4)

# =========================================================
# 5. 백스톱 인덱스 + 데드락 주의
# =========================================================
s = slide()
bg(s)
header(s, "동시성 #3 · 안전망", "DB 레벨 불변식 + 데드락 관점")
code_box(s, Inches(0.6), Inches(1.6), Inches(7.4), Inches(2.4), [
    "-- 같은 쌍에 active 매칭이 둘 이상 생기지 못하게 하는",
    "-- 부분 유니크 인덱스 (status='active' 일 때만)",
    "create unique index matches_active_pair_uniq",
    "  on matches (least(user_a,user_b),",
    "              greatest(user_a,user_b))",
    "  where status = 'active';",
], fontsize=12.5)
card(s, Inches(0.6), Inches(4.2), Inches(5.95), Inches(2.3), "왜 인덱스까지?",
     ["락은 '코드가 항상 락을 잡을 때'만 유효",
      "인덱스 제약은 경로와 무관하게 항상 성립",
      "→ 다중 방어(defense in depth)",
      "OS의 '불변식(invariant) 보장'과 같은 발상"])
card(s, Inches(6.75), Inches(1.6), Inches(5.95), Inches(2.0), "교착 상태(Deadlock) 주의",
     ["여러 락을 잡으면 순환 대기 → 데드락",
      "여기선 단일 키('match_creation') 사용",
      "→ 락 순서 단순화로 데드락 회피"], accent=INK)
card(s, Inches(6.75), Inches(3.75), Inches(5.95), Inches(2.75), "재매칭 함수도 같은 키",
     ["find_or_create_match 와 find_additional_match 가",
      "동일 키를 공유 → 두 함수끼리도 직렬화",
      "교과서의 '모든 임계구역이 같은 락 사용' 원칙",
      "= 상호배제가 전역적으로 성립"], accent=CORAL)
footer(s, 5)

# =========================================================
# 6. 캐시
# =========================================================
s = slide()
bg(s)
header(s, "메모리 계층 #1", "캐싱(Caching) — 멜론 검색/차트 인메모리 캐시")
code_box(s, Inches(0.6), Inches(1.6), Inches(7.4), Inches(3.5), [
    "const store = new Map<string, Entry>();  // 서버 메모리",
    "",
    "function getCache(key){",
    "  const e = store.get(key);",
    "  if(!e) return null;                 // miss",
    "  if(Date.now() > e.expiresAt){       // 만료(eviction)",
    "    store.delete(key); return null;",
    "  }",
    "  return e.data;                      // hit",
    "}",
    "",
    "CACHE_TTL = { chart: 5분, search: 10분,",
    "              searchFail: 30초 }  // 장애 응답은 짧게",
], fontsize=12.5)
card(s, Inches(8.25), Inches(1.6), Inches(4.45), Inches(2.35), "OS 매핑",
     ["메모리 캐시 = CPU 캐시/페이지 캐시",
      "hit/miss, TTL 만료 = 캐시 교체 정책",
      "같은 검색어 반복 = 시간적 지역성",
      "외부 API(느린 디스크) → 메모리(빠름)"], accent=CORAL)
card(s, Inches(8.25), Inches(4.15), Inches(4.45), Inches(2.35), "왜 TTL을 다르게?",
     ["정상 결과는 길게(10분) 캐시 → 부하↓",
      "장애 폴백은 30초만 → 빠른 복구",
      "= 캐시 일관성 vs 성능 trade-off",
      "신선도(freshness)와 비용의 균형"], accent=INK)
footer(s, 6)

# =========================================================
# 7. 캐시 효과 (개념도)
# =========================================================
s = slide()
bg(s)
header(s, "메모리 계층 #2", "지역성과 응답 흐름 — 캐시가 막아주는 것")
# 흐름 박스
def flowbox(x, title, sub, color, tcol=WHITE):
    rect(s, x, Inches(2.4), Inches(2.5), Inches(1.5), color)
    txt(s, x, Inches(2.7), Inches(2.5), Inches(0.6), [[(title, 16, True, tcol)]], align=PP_ALIGN.CENTER)
    txt(s, x, Inches(3.25), Inches(2.5), Inches(0.6), [[(sub, 11.5, False, tcol)]], align=PP_ALIGN.CENTER)

flowbox(Inches(0.6), "검색 요청", "사용자 입력", LIGHT, INK)
txt(s, Inches(3.15), Inches(2.9), Inches(0.7), Inches(0.5), [[("→", 28, True, CORAL)]], align=PP_ALIGN.CENTER)
flowbox(Inches(3.85), "캐시 조회", "Map.get(key)", CORAL_SOFT, INK)
txt(s, Inches(6.4), Inches(2.9), Inches(0.7), Inches(0.5), [[("→", 28, True, CORAL)]], align=PP_ALIGN.CENTER)
flowbox(Inches(7.1), "HIT", "즉시 응답 ⚡", CORAL)
flowbox(Inches(10.0), "MISS", "멜론 API 호출", INK)
txt(s, Inches(7.1), Inches(4.0), Inches(2.5), Inches(0.4), [[("빠름 (메모리)", 11, True, CORAL)]], align=PP_ALIGN.CENTER)
txt(s, Inches(10.0), Inches(4.0), Inches(2.5), Inches(0.4), [[("느림 (외부 네트워크)", 11, True, GRAY)]], align=PP_ALIGN.CENTER)
rect(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.7), LIGHT)
txt(s, Inches(0.85), Inches(5.05), Inches(11.6), Inches(0.5), [[("핵심 통찰", 15, True, CORAL)]])
txt(s, Inches(0.85), Inches(5.5), Inches(11.6), Inches(1.0),
    [[("• 느린 계층(외부 API) 앞에 빠른 계층(메모리)을 두는 것 = 운영체제 메모리 계층구조의 원리.", 13.5, False, INK)],
     [("• 인기 검색어가 반복(지역성)될수록 hit율↑ → 평균 응답시간과 외부 부하가 동시에 감소.", 13.5, False, INK)]],
    space_after=6)
footer(s, 7)

# =========================================================
# 8. 프로세스 상태 / 스케줄링 / 큐
# =========================================================
s = slide()
bg(s)
header(s, "스케줄링", "프로세스 상태 & 큐 — 가입 승인 / 매칭 / '한 명 더' 대기")
states = [
    ("가입 → 승인 대기", "운영팀 검수 큐에서 대기", "New / Ready"),
    ("승인 완료(active)", "매칭·채팅 자원 사용 가능", "Running"),
    ("'한 명 더' 대기", "즉시 연결 X, 확인 후 처리", "Waiting/Blocked"),
    ("탈퇴 / 매칭 종료", "자원 회수, 상태 종료", "Terminated"),
]
x = Inches(0.6)
for t, d, os_ in states:
    rect(s, x, Inches(1.7), Inches(2.95), Inches(2.2), LIGHT)
    rect(s, x, Inches(1.7), Inches(2.95), Inches(0.5), CORAL)
    txt(s, x, Inches(1.78), Inches(2.95), Inches(0.4), [[(os_, 13, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), Inches(2.35), Inches(2.65), Inches(0.6), [[(t, 13.5, True, INK)]])
    txt(s, x + Inches(0.15), Inches(2.95), Inches(2.65), Inches(0.9), [[(d, 12, False, GRAY)]])
    x = x + Inches(3.05)
rect(s, Inches(0.6), Inches(4.25), Inches(12.1), Inches(2.35), CORAL_SOFT)
txt(s, Inches(0.85), Inches(4.4), Inches(11.6), Inches(0.5), [[("📖 OS 매핑", 15, True, CORAL)]])
txt(s, Inches(0.85), Inches(4.85), Inches(11.6), Inches(1.6),
    [[("• 승인 대기 = 작업 큐(Job Queue)에서 스케줄러(운영팀)의 선택을 기다리는 New/Ready 상태.", 13.5, False, INK)],
     [("• '한 명 더'를 즉시 연결하지 않고 대기시키는 것 = 자원(매칭 슬롯) 경합을 줄이는 admission control.", 13.5, False, INK)],
     [("• 미승인 사용자의 채팅 차단 = 자원 접근 전 상태 검사 → 잘못된 상태 전이를 막는 가드.", 13.5, False, INK)]],
    space_after=8)
footer(s, 8)

# =========================================================
# 9. IPC / Realtime
# =========================================================
s = slide()
bg(s)
header(s, "프로세스 간 통신", "IPC — 실시간 채팅은 메시지 패싱(Message Passing)")
card(s, Inches(0.6), Inches(1.65), Inches(5.95), Inches(2.4), "앱에서 일어나는 일",
     ["A가 메시지 전송 → DB messages INSERT",
      "Supabase Realtime이 변경을 구독자에게 push",
      "B의 화면에 즉시 반영 (수정/삭제/열람도 전파)",
      "= 두 클라이언트가 채널로 메시지를 주고받음"])
card(s, Inches(6.75), Inches(1.65), Inches(5.95), Inches(2.4), "OS 매핑",
     ["서로 다른 프로세스(브라우저 2개)의 통신",
      "공유 메모리가 아닌 '메시지 전달' 모델",
      "DB/Realtime 채널 = 커널이 중개하는 메일박스",
      "구독(subscribe) = 채널에 receive 등록"], accent=INK)
rect(s, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.3), LIGHT)
txt(s, Inches(0.85), Inches(4.45), Inches(11.6), Inches(0.5),
    [[("동기화 포인트", 15, True, CORAL)]])
txt(s, Inches(0.85), Inches(4.9), Inches(11.6), Inches(1.5),
    [[("• 'replica identity' 설정으로 UPDATE 전체 행을 전파 → 양쪽 상태를 일관되게 유지(consistency).", 13.5, False, INK)],
     [("• 메시지 패싱은 공유 변수 락이 필요 없어 경쟁 상태가 적다 → IPC 방식 선택의 trade-off.", 13.5, False, INK)],
     [("• 비교: 매칭(슬라이드 3~5)은 공유 메모리(공유 테이블) 모델이라 반대로 락이 필요했다.", 13.5, False, INK)]],
    space_after=6)
footer(s, 9)

# =========================================================
# 10. 자원 회수 / 파일 수명
# =========================================================
s = slide()
bg(s)
header(s, "자원 관리", "자원 수명 & 회수 — '한 번만 보기' 사진 / '펑' 메시지")
code_box(s, Inches(0.6), Inches(1.6), Inches(6.6), Inches(3.0), [
    "-- messages 상태 컬럼",
    "viewed_at  : 수신자 열람 시각",
    "view_once  : 1회 열람 후 스토리지 원본 삭제",
    "disappear  : '펑' — 읽으면 자동 소멸(soft delete)",
    "deleted_at : 소멸/삭제 시각 (null=살아있음)",
], fontsize=13)
card(s, Inches(7.4), Inches(1.6), Inches(5.3), Inches(3.0), "수명 주기",
     ["생성 → 열람(viewed_at) → 회수(원본 delete)",
      "한 번 쓰고 버리는 자원 = 명시적 해제",
      "참조가 끝나면 즉시 자원을 반납하는 패턴"], accent=CORAL)
rect(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.8), CORAL_SOFT)
txt(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(0.5), [[("📖 OS 매핑", 15, True, CORAL)]])
txt(s, Inches(0.85), Inches(5.4), Inches(11.6), Inches(1.1),
    [[("• view_once 사진의 '열람 후 원본 삭제' = 사용 끝난 자원의 즉시 회수(메모리 해제/파일 unlink와 동형).", 13.5, False, INK)],
     [("• 'viewed_at에 따라 소멸' = 조건 충족 시 상태 전이 → 동기화(열람 이벤트)와 자원 해제를 묶은 설계.", 13.5, False, INK)]],
    space_after=6)
footer(s, 10)

# =========================================================
# 11. 보호 / 접근제어 / 이중 모드
# =========================================================
s = slide()
bg(s)
header(s, "보호와 보안", "보호(Protection) & 접근제어 — RLS / service_role / 서버 전용")
card(s, Inches(0.6), Inches(1.65), Inches(3.85), Inches(2.5), "사용자 모드",
     ["브라우저의 anon/로그인 키",
      "RLS 정책이 허용한 행만 접근",
      "'본인 메시지만 수정/삭제' 정책",
      "= 최소 권한(least privilege)"])
card(s, Inches(4.65), Inches(1.65), Inches(3.85), Inches(2.5), "커널(특권) 모드",
     ["서버 전용 service_role 키",
      "RLS 우회 가능한 강한 권한",
      "절대 클라이언트에 노출 금지",
      "= 특권 명령은 커널만"], accent=INK)
card(s, Inches(8.7), Inches(1.65), Inches(4.0), Inches(2.5), "경계(boundary)",
     [".server.ts = 서버에서만 실행",
      "DB 쿼리는 repos/*.server 안에서만",
      "클라이언트↔서버 권한 분리",
      "= 사용자/커널 모드 전환"], accent=CORAL)
rect(s, Inches(0.6), Inches(4.4), Inches(12.1), Inches(2.2), LIGHT)
txt(s, Inches(0.85), Inches(4.55), Inches(11.6), Inches(0.5), [[("핵심", 15, True, CORAL)]])
txt(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.4),
    [[("• RLS(Row Level Security) = OS의 접근제어 행렬/권한 비트 — '누가 어떤 자원에 무엇을 할 수 있나'.", 13.5, False, INK)],
     [("• service_role과 anon 키의 분리 = 이중 모드(user/kernel)와 권한 도메인 분리의 데이터베이스 버전.", 13.5, False, INK)],
     [("• 색상 가드(pre-commit) · 권한 허용목록 = 정책을 자동으로 강제하는 reference monitor의 축소판.", 13.5, False, INK)]],
    space_after=6)
footer(s, 11)

# =========================================================
# 12. 정리
# =========================================================
s = slide()
bg(s, INK)
rect(s, 0, 0, Inches(0.25), SH, CORAL)
txt(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.8),
    [[("정리 — 시험에 이렇게 답하면 된다", 32, True, WHITE)]])
summ = [
    ("임계구역/상호배제", "매칭의 SELECT~INSERT를 advisory lock으로 직렬화 + 유니크 인덱스로 불변식 보장"),
    ("경쟁 상태/이중검사", "락 획득 후 재확인(double-checked) — 대기 중 상태가 바뀔 수 있으니"),
    ("캐싱/지역성", "느린 외부 API 앞에 메모리 캐시(TTL) → hit/miss·교체 = 메모리 계층구조"),
    ("프로세스 상태/스케줄링", "가입 승인 큐·'한 명 더' 대기 = Ready/Waiting 상태와 admission control"),
    ("IPC", "실시간 채팅 = 메시지 패싱(공유 메모리 아님) → 락 없이 통신"),
    ("자원 회수", "한 번만 보기/펑 = 사용 후 즉시 해제(unlink) + 상태 전이 동기화"),
    ("보호/이중 모드", "RLS·service_role 분리 = 접근제어 + user/kernel 권한 도메인"),
]
y = Inches(1.7)
for k, v in summ:
    rect(s, Inches(0.9), y, Inches(3.5), Inches(0.62), CORAL)
    txt(s, Inches(1.0), y + Inches(0.13), Inches(3.3), Inches(0.4), [[(k, 13, True, WHITE)]])
    txt(s, Inches(4.6), y + Inches(0.08), Inches(8.2), Inches(0.6),
        [[(v, 12.5, False, RGBColor(0xE0, 0xE0, 0xE8))]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.72)

prs.save("/home/user/music-dating-app-keonuu/뮤직매치_운영체제_개념정리.pptx")
print("saved")
